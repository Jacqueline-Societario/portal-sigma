"""
blueprints/informativos.py — Gerador de Informativos (PPTX + PDF)
Layout baseado nos moldes de referência Sigma:
  - Capa: painel esquerdo #A72C31 + painel direito branco (ISSNET style)
  - Slides de passo: header fino #A72C31 + barra azul + blocos (NF-SP style)
  - Slide de contato final: ISSNET style exato
PPTX e PDF usam as mesmas medidas — saída idêntica.
"""

import os, uuid, json, io, time, textwrap
from flask import (Blueprint, render_template, request,
                   redirect, url_for, send_file, current_app)
from blueprints.auth import login_obrigatorio
from werkzeug.utils import secure_filename

# ── python-pptx ───────────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── ReportLab (canvas drawing) ────────────────────────────────────────────────
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor, white, black
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.utils import ImageReader

# ── PIL ────────────────────────────────────────────────────────────────────────
from PIL import Image as PILImage

informativos_bp = Blueprint('informativos', __name__, url_prefix='/informativos')

ALLOWED_IMG = {'png', 'jpg', 'jpeg', 'gif', 'webp'}

# ── Cache de documentos gerados ───────────────────────────────────────────────
_CACHE: dict = {}

# ── Dados fixos Sigma ─────────────────────────────────────────────────────────
_SIGMA = {
    'tel1':    '(62) 3095-6240',
    'tel1_d':  'Apenas Ligação',
    'tel2':    '(62) 9 9699-9939',
    'tel2_d':  'Apenas WhatsApp',
    'email1':  'atendimento@gsigma.com.br',
    'email2':  'contato@gsigma.com.br',
    'email3':  'societario@gsigma.com.br',
    'site':    'www.gsigma.com.br',
    'end1':    'Av. T-2, nº 471 — Sala 507',
    'end2':    'Focus Business Center, Setor Bueno',
    'end3':    'Goiânia/GO  •  CEP 74210-005',
    'horario': 'Seg–Sex 8h–17h45',
}

# ── Cores ─────────────────────────────────────────────────────────────────────
_RED   = RGBColor(0xA7, 0x2C, 0x31)
_DKRED = RGBColor(0x80, 0x1E, 0x24)
_BLUE  = RGBColor(0x1A, 0x3A, 0x6E)
_LBLUE = RGBColor(0x13, 0x51, 0xB4)
_YELL  = RGBColor(0xFF, 0xCD, 0x07)
_GRN   = RGBColor(0x16, 0x88, 0x21)
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_DARK  = RGBColor(0x1A, 0x1A, 0x1A)
_GRAY  = RGBColor(0x55, 0x55, 0x55)
_LGRAY = RGBColor(0xF5, 0xF5, 0xF5)
_PINK  = RGBColor(0xFD, 0xF0, 0xF0)

# ── Constantes de layout (em inches) — medidas dos moldes ─────────────────────
_SW   = 13.33   # slide width
_SH   = 7.50    # slide height

# Header (NF-SP style — slides de passo)
_H_H  = 0.55    # header height
_H_Y  = 0.0     # header top

# Footer (NF-SP style)
_F_Y  = 7.12    # footer top
_F_H  = 0.38    # footer height

# Step bar (NF-SP style)
_B_Y  = 0.65    # bar top
_B_H  = 0.65    # bar height
_B_BW = 0.50    # badge width/height

# Content area (after step bar)
_C_X  = 0.45    # content left margin
_C_W  = 12.43   # content width
_C_Y0 = 1.38    # first content block top

# Block heights (NF-SP measurements)
_LBL_H  = 0.28  # label bar height
_LBL_TY = 0.04  # text offset inside label
_CONT_H  = 1.00  # main content box height
_LINK_H  = 0.72  # clickable box height
_GRNH    = 2.19  # green "o que verá" box height
_WARN_H  = 0.52  # warning box height

# Contact slide (ISSNET style)
_K_H_H  = 1.20  # contact header height
_K_DSC_Y = 1.45  # description box top
_K_DSC_H = 1.20  # description box height
_K_CRD_Y = 2.90  # cards top
_K_CRD_H = 2.45  # cards height
_K_CRD_W = 2.90  # card width
_K_FT_Y  = 5.55  # footer top
_K_FT_H  = 1.95  # footer height
_K_CARDS_X = [0.50, 3.72, 6.94, 10.16]  # card left positions

# ─────────────────────────────────────────────────────────────────────────────
# PPTX HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def _i(v): return Inches(v)

def _rect(slide, l, t, w, h, fill_rgb):
    """Adiciona retângulo sólido ao slide."""
    shp = slide.shapes.add_shape(1, _i(l), _i(t), _i(w), _i(h))
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    return shp

def _txt(slide, l, t, w, h, text, size, bold=False, italic=False,
         color=None, align=PP_ALIGN.LEFT, wrap=True):
    """Adiciona caixa de texto ao slide."""
    if color is None:
        color = _DARK
    txb = slide.shapes.add_textbox(_i(l), _i(t), _i(w), _i(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def _img_pptx(slide, l, t, w, h, pil_img):
    """Insere imagem PIL no slide."""
    buf = io.BytesIO()
    rgb = pil_img.convert('RGB')
    rgb.save(buf, format='JPEG', quality=85)
    buf.seek(0)
    slide.shapes.add_picture(buf, _i(l), _i(t), _i(w), _i(h))

def _label_block(slide, x, y, w, h, fill_rgb, text, txt_color=None):
    """Barra de rótulo + texto (O QUE FAZER, CLIQUE, etc.)."""
    if txt_color is None:
        txt_color = _WHITE
    _rect(slide, x, y, w, h, fill_rgb)
    _txt(slide, x + 0.12, y + _LBL_TY, w - 0.20, h, text,
         size=8, bold=True, color=txt_color)

def _content_block(slide, x, y, w, h, fill_rgb, text, size=11, color=None):
    """Bloco de conteúdo com fundo colorido."""
    if color is None:
        color = _DARK
    _rect(slide, x, y, w, h, fill_rgb)
    _txt(slide, x + 0.15, y + 0.10, w - 0.20, h - 0.15,
         text, size=size, color=color)

# ─────────────────────────────────────────────────────────────────────────────
# PPTX — CABEÇALHO E RODAPÉ PADRÃO (NF-SP style)
# ─────────────────────────────────────────────────────────────────────────────

def _pptx_header(slide, guide_title=''):
    """Header fino #A72C31 com 'SIGMA CONTABILIDADE' + título do guia."""
    _rect(slide, 0, 0, _SW, _H_H, _RED)
    _txt(slide, 0.45, 0.08, 5.5, 0.40, 'SIGMA CONTABILIDADE',
         size=10, bold=True, color=_WHITE)
    if guide_title:
        _txt(slide, 6.0, 0.08, 6.88, 0.40, guide_title,
             size=10, color=_WHITE, align=PP_ALIGN.RIGHT)

def _pptx_footer(slide, page_num, total, footer_text=''):
    """Rodapé #F5F5F5 com texto e número de página."""
    _rect(slide, 0, _F_Y, _SW, _F_H, _LGRAY)
    if not footer_text:
        footer_text = (f'Sigma Contabilidade  •  Goiânia – GO  •  '
                       f'Dúvidas? Entre em contato: {_SIGMA["tel2"]}')
    _txt(slide, 0.45, 7.16, 11.20, 0.30, footer_text,
         size=8, color=_GRAY)
    _txt(slide, 12.13, 7.16, 0.80, 0.30,
         f'{page_num}/{total}', size=8, color=_GRAY,
         align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# PPTX — SLIDE DE CAPA (ISSNET style)
# ─────────────────────────────────────────────────────────────────────────────

def _slide_capa(prs, dados):
    """
    Capa estilo ISSNET:
    - Painel esquerdo #A72C31: SIGMA grande + tagline
    - Painel direito branco: subtítulo + título grande + cliente
    """
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Painel esquerdo
    _rect(sl, 0.0, 0.0, 5.20, _SH, _RED)
    _rect(sl, 5.20, 0.0, 0.06, _SH, _DKRED)  # divisor
    # Painel direito
    _rect(sl, 5.26, 0.0, 8.07, _SH, _WHITE)

    # Esquerda: SIGMA + CONTABILIDADE + tagline
    _txt(sl, 0.55, 1.60, 4.20, 1.00, 'SIGMA',
         size=54, bold=True, color=_WHITE)
    _txt(sl, 0.55, 2.55, 4.20, 0.50, 'CONTABILIDADE',
         size=18, color=RGBColor(0xFF, 0xCC, 0xCC))
    _txt(sl, 0.55, 3.05, 4.20, 0.40, 'Além da Contabilidade',
         size=12, italic=True, color=RGBColor(0xFF, 0xDD, 0xDD))
    # Linha decorativa
    _rect(sl, 0.55, 3.55, 3.50, 0.04, RGBColor(0xFF, 0xCC, 0xCC))

    # Direita: subtítulo + título + cliente/descrição
    subtitulo = dados.get('subtitulo', '') or 'Guia passo a passo — simples e rápido'
    titulo    = dados.get('titulo', 'Informativo Sigma')
    cliente   = dados.get('cliente', '')

    _txt(sl, 5.50, 1.50, 7.40, 0.50, subtitulo,
         size=14, italic=True, color=_RED)
    _txt(sl, 5.50, 2.10, 7.40, 2.20, titulo,
         size=36, bold=True, color=_DARK)

    if cliente:
        _txt(sl, 5.50, 4.50, 7.40, 0.50, f'Preparado para: {cliente}',
             size=13, color=_GRAY)

    # Rodapé esquerdo
    _txt(sl, 0.55, 6.90, 4.20, 0.40,
         f'{_SIGMA["site"]}  •  {_SIGMA["tel1"]}',
         size=10, italic=True, color=RGBColor(0xAA, 0xAA, 0xAA))

# ─────────────────────────────────────────────────────────────────────────────
# PPTX — SLIDE DE AVISOS (slide "Antes de Começar")
# ─────────────────────────────────────────────────────────────────────────────

def _slide_avisos(prs, avisos, num_slide, total, titulo_guia):
    """Slide com lista de avisos importantes."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    _pptx_header(sl, titulo_guia)
    _pptx_footer(sl, num_slide, total)

    # Barra de título
    _rect(sl, 0, _B_Y, _SW, _B_H, _BLUE)
    _txt(sl, 0.45, 0.75, 12.43, 0.45,
         '⚠️  ANTES DE COMEÇAR — INFORMAÇÕES IMPORTANTES',
         size=16, bold=True, color=_WHITE)

    # Caixas de aviso (até 6 avisos em 2 colunas de 3)
    y_start = 1.42
    col_w   = 6.00
    col_gap = 0.43
    box_h   = 1.50
    row_gap = 0.08

    fills = [
        RGBColor(0xFF, 0xF8, 0xDC),  # amarelo suave
        RGBColor(0xDB, 0xE8, 0xFB),  # azul suave
        RGBColor(0xE3, 0xF5, 0xE1),  # verde suave
        RGBColor(0xFF, 0xF8, 0xDC),
        RGBColor(0xDB, 0xE8, 0xFB),
        RGBColor(0xE3, 0xF5, 0xE1),
    ]
    txt_colors = [_RED, _BLUE, _GRN, _RED, _BLUE, _GRN]

    for i, aviso in enumerate(avisos[:6]):
        col = i % 2
        row = i // 2
        x = _C_X + col * (col_w + col_gap)
        y = y_start + row * (box_h + row_gap)
        _rect(sl, x, y, col_w, box_h, fills[i])
        _txt(sl, x + 0.15, y + 0.15, col_w - 0.25, box_h - 0.25,
             aviso, size=13, bold=True, color=txt_colors[i])

# ─────────────────────────────────────────────────────────────────────────────
# PPTX — SLIDE DE PASSO (NF-SP style)
# ─────────────────────────────────────────────────────────────────────────────

def _slide_passo(prs, passo, num_slide, total, titulo_guia, imagem=None):
    """
    Slide de passo estilo NF-SP:
    - Header fino + rodapé
    - Barra azul com badge amarelo + título
    - Bloco "O QUE FAZER" (vermelho) + conteúdo cinza
    - Se link: bloco "CLIQUE / TOQUE AQUI" (azul) + URL
    - Se imagem: painel direito com imagem
    """
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    _pptx_header(sl, titulo_guia)
    _pptx_footer(sl, num_slide, total)

    # Barra de passo (azul escuro)
    _rect(sl, 0, _B_Y, _SW, _B_H, _BLUE)
    # Badge amarelo com número
    _rect(sl, _C_X, 0.725, _B_BW, _B_BW, _YELL)
    _txt(sl, _C_X, 0.755, _B_BW, 0.46,
         str(passo.get('num', '')),
         size=20, bold=True, color=_BLUE, align=PP_ALIGN.CENTER)
    # Título do passo
    _txt(sl, 1.13, 0.79, 11.73, 0.46,
         passo.get('titulo', ''), size=17, bold=True, color=_WHITE)

    # Determina se tem imagem (painel à direita)
    has_img = imagem is not None
    c_w = 6.90 if has_img else _C_W
    c_x = _C_X

    # Bloco "O QUE FAZER"
    y = _C_Y0
    _label_block(sl, c_x, y, c_w, _LBL_H, _RED, 'O QUE FAZER')

    # Conteúdo principal
    y += _LBL_H
    descricao = passo.get('descricao', '')
    _content_block(sl, c_x, y, c_w, _CONT_H,
                   RGBColor(0xF5, 0xF5, 0xF5), descricao, size=11)

    # Bloco "CLIQUE / TOQUE AQUI" (sempre incluso — instrução genérica se sem link)
    y += _CONT_H
    link = passo.get('link', '')
    _label_block(sl, c_x, y, c_w, _LBL_H, _LBLUE,
                 '👉  CLIQUE / TOQUE AQUI')
    y += _LBL_H

    link_text = link if link else 'Siga as instruções na tela conforme descrito acima.'
    _content_block(sl, c_x, y, c_w, _LINK_H,
                   RGBColor(0xDB, 0xE8, 0xFB), link_text, size=11)

    # Bloco "O QUE VOCÊ VERÁ" (verde) — espaço restante
    y += _LINK_H
    resultado = passo.get('resultado', '')
    available_h = _F_Y - y - 0.08
    if available_h > 0.4:
        _label_block(sl, c_x, y, c_w, _LBL_H, _GRN,
                     '🖥️  O QUE VOCÊ VERÁ NA TELA')
        y += _LBL_H
        grn_h = min(available_h - _LBL_H, _GRNH)
        if grn_h > 0.2:
            txt_resultado = resultado if resultado else 'Confira o resultado na tela e prossiga para o próximo passo.'
            _content_block(sl, c_x, y, c_w, grn_h,
                           RGBColor(0xE3, 0xF5, 0xE1), txt_resultado, size=11)

    # Painel direito com imagem
    if has_img:
        _rect(sl, 7.66, _C_Y0, 5.22, 5.62, _LGRAY)
        # Centralizar imagem dentro do painel
        panel_x, panel_y = 7.76, _C_Y0 + 0.15
        panel_w, panel_h = 5.02, 5.32
        # Calcular proporção para caber no painel
        img_w, img_h = imagem.size
        ratio = min(panel_w / (img_w / 96), panel_h / (img_h / 96))
        draw_w = (img_w / 96) * ratio
        draw_h = (img_h / 96) * ratio
        off_x = panel_x + (panel_w - draw_w) / 2
        off_y = panel_y + (panel_h - draw_h) / 2
        _img_pptx(sl, off_x, off_y, draw_w, draw_h, imagem)

# ─────────────────────────────────────────────────────────────────────────────
# PPTX — SLIDE DE LINKS
# ─────────────────────────────────────────────────────────────────────────────

def _slide_links(prs, links, num_slide, total, titulo_guia):
    """Slide com lista de links e recursos adicionais."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    _pptx_header(sl, titulo_guia)
    _pptx_footer(sl, num_slide, total)

    # Barra
    _rect(sl, 0, _B_Y, _SW, _B_H, _BLUE)
    _txt(sl, 0.45, 0.75, 12.43, 0.45,
         '🔗  LINKS E RECURSOS ADICIONAIS',
         size=16, bold=True, color=_WHITE)

    # Cabeçalho da tabela
    y = _C_Y0
    _rect(sl, _C_X, y, _C_W, 0.35, _RED)
    _txt(sl, _C_X + 0.15, y + 0.05, 4.00, 0.25,
         'DESCRIÇÃO', size=9, bold=True, color=_WHITE)
    _txt(sl, _C_X + 4.30, y + 0.05, 8.00, 0.25,
         'ENDEREÇO / LINK', size=9, bold=True, color=_WHITE)

    y += 0.35
    row_h = 0.55
    for i, lnk in enumerate(links[:9]):
        fill = RGBColor(0xF5, 0xF5, 0xF5) if i % 2 == 0 else _WHITE
        _rect(sl, _C_X, y, _C_W, row_h, fill)
        _txt(sl, _C_X + 0.15, y + 0.08, 4.00, row_h - 0.10,
             lnk.get('label', ''), size=11, color=_DARK)
        _txt(sl, _C_X + 4.30, y + 0.08, 8.00, row_h - 0.10,
             lnk.get('url', ''), size=10, color=_LBLUE)
        y += row_h

# ─────────────────────────────────────────────────────────────────────────────
# PPTX — SLIDE DE CONTATO FINAL (ISSNET style — medidas exatas do molde)
# ─────────────────────────────────────────────────────────────────────────────

def _slide_contato_final(prs, num_slide, total, titulo_guia):
    """
    Slide de contato — IDÊNTICO ao molde ISSNET Slide 7:
    - Header alto #A72C31 (1.20")
    - Caixa de descrição #FDF0F0
    - 4 cards #F5F5F5 com ícone + rótulo + conteúdo
    - Rodapé #A72C31 com nome Sigma + tagline
    """
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Header (ISSNET — mais alto que os outros slides) ──
    _rect(sl, 0.0, 0.0, _SW, _K_H_H, _RED)
    _rect(sl, 0.0, _K_H_H, _SW, 0.06, _DKRED)
    _txt(sl, 0.60, 0.28, 12.00, 0.65,
         'PRECISA DE AJUDA?', size=24, bold=True, color=_WHITE)
    _txt(sl, 0.60, 0.82, 12.00, 0.35,
         'Nossa equipe está pronta para te ajudar',
         size=13, italic=True, color=RGBColor(0xFF, 0xDD, 0xDD))

    # ── Caixa de descrição ──
    _rect(sl, 0.50, _K_DSC_Y, 12.30, _K_DSC_H, _PINK)
    _txt(sl, 0.75, 1.58, 11.80, 0.95,
         (f'Entre em contato com a Sigma Contabilidade para tirar dúvidas, '
          f'solicitar suporte ou agendar uma conversa com nossos especialistas. '
          f'Estamos disponíveis de {_SIGMA["horario"]}.'),
         size=14, color=_DARK)

    # ── 4 Cards ──
    card_data = [
        ('📞', 'Telefones',
         f'{_SIGMA["tel1"]}\n{_SIGMA["tel1_d"]}\n\n'
         f'{_SIGMA["tel2"]}\n{_SIGMA["tel2_d"]}'),
        ('✉️', 'E-mail',
         f'{_SIGMA["email1"]}\n{_SIGMA["email2"]}\n{_SIGMA["email3"]}'),
        ('🌐', 'Site',
         f'{_SIGMA["site"]}\n\nAcesse e conheça\nnossas soluções'),
        ('📍', 'Endereço',
         f'{_SIGMA["end1"]}\n{_SIGMA["end2"]}\n{_SIGMA["end3"]}'),
    ]

    for i, (icon, label, content) in enumerate(card_data):
        cx = _K_CARDS_X[i]
        _rect(sl, cx, _K_CRD_Y, _K_CRD_W, _K_CRD_H, _LGRAY)
        # Ícone
        _txt(sl, cx + 0.10, _K_CRD_Y + 0.10, 0.60, 0.60,
             icon, size=26, align=PP_ALIGN.CENTER)
        # Rótulo
        _txt(sl, cx + 0.80, _K_CRD_Y + 0.15, 1.80, 0.40,
             label, size=13, bold=True, color=_RED)
        # Conteúdo
        _txt(sl, cx + 0.15, _K_CRD_Y + 0.62, _K_CRD_W - 0.20, 1.70,
             content, size=12, color=_GRAY)

    # ── Rodapé bordeaux (ISSNET style) ──
    _rect(sl, 0.0, _K_FT_Y, _SW, _K_FT_H, _RED)
    _txt(sl, 0.60, 5.75, 8.00, 0.65,
         'SIGMA CONTABILIDADE', size=26, bold=True, color=_WHITE)
    _txt(sl, 0.60, 6.35, 7.00, 0.40,
         'Além da Contabilidade', size=14, italic=True,
         color=RGBColor(0xFF, 0xCC, 0xCC))
    _txt(sl, 0.60, 6.85, 12.00, 0.35,
         f'{_SIGMA["site"]}  •  {_SIGMA["tel1"]}  •  {_SIGMA["tel2"]}  •  {_SIGMA["email1"]}',
         size=10, color=RGBColor(0xFF, 0xCC, 0xCC))
    _txt(sl, 12.80, 7.20, 0.45, 0.25,
         f'{num_slide}/{total}', size=9, color=RGBColor(0xAA, 0xAA, 0xAA),
         align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────────────
# PPTX — GERADOR PRINCIPAL
# ─────────────────────────────────────────────────────────────────────────────

def _gerar_pptx(dados, path_out):
    """Gera o PPTX completo e salva em path_out."""
    prs = Presentation()
    prs.slide_width  = Inches(_SW)
    prs.slide_height = Inches(_SH)

    titulo_guia = dados.get('titulo', 'Informativo Sigma')
    passos      = dados.get('passos', [])
    avisos      = dados.get('avisos', [])
    links       = dados.get('links', [])
    imagens     = dados.get('imagens', [])

    # Contar total de slides
    total = 1  # capa
    if avisos:
        total += 1
    total += len(passos)
    if links:
        total += 1
    total += 1  # contato

    num = 1

    # 1. Capa
    _slide_capa(prs, dados)
    num += 1

    # 2. Avisos (se houver)
    if avisos:
        _slide_avisos(prs, avisos, num, total, titulo_guia)
        num += 1

    # 3. Passos
    for i, passo in enumerate(passos):
        img = imagens[i] if i < len(imagens) else None
        _slide_passo(prs, passo, num, total, titulo_guia, img)
        num += 1

    # 4. Links (se houver)
    if links:
        _slide_links(prs, links, num, total, titulo_guia)
        num += 1

    # 5. Contato final
    _slide_contato_final(prs, num, total, titulo_guia)

    prs.save(path_out)

# ─────────────────────────────────────────────────────────────────────────────
# PDF — CANVAS DRAWING (medidas idênticas ao PPTX)
# Usa página 13.33" × 7.5" = 960pt × 540pt (72pt/inch)
# ─────────────────────────────────────────────────────────────────────────────

# Página com mesmas proporções do PPTX
_PDF_W = 13.33 * 72   # 959.76pt
_PDF_H = 7.50  * 72   # 540pt

def _pi(v): return v * 72   # inches → points

def _pdf_rect(c, l, t, w, h, fill_hex):
    """Retângulo sólido no canvas PDF. t = top em inches (convertido para RL coords)."""
    c.setFillColor(HexColor(fill_hex))
    c.setStrokeColor(HexColor(fill_hex))
    # ReportLab: y=0 embaixo. Converter: rl_y = height - top - h
    x = _pi(l)
    y = _PDF_H - _pi(t) - _pi(h)
    c.rect(x, y, _pi(w), _pi(h), fill=1, stroke=0)

def _pdf_txt(c, l, t, w, h, text, size, bold=False, italic=False,
             color='#1A1A1A', align='L', v_center=False):
    """Texto no canvas PDF."""
    c.setFillColor(HexColor(color))
    c.setFont('Helvetica-Bold' if bold else ('Helvetica-Oblique' if italic else 'Helvetica'), size)

    x = _pi(l)
    # Posição y: topo da caixa menos o tamanho da fonte (baseline)
    y = _PDF_H - _pi(t) - size
    if v_center:
        y = _PDF_H - _pi(t) - _pi(h) / 2 - size / 2

    max_w = _pi(w)
    # Quebrar texto se necessário
    lines = _pdf_wrap(c, text, size, bold, italic, max_w)
    line_h = size * 1.4
    for i, line in enumerate(lines):
        ly = y - i * line_h
        if ly < _PDF_H - _pi(t) - _pi(h):
            break
        if align == 'R':
            tw = c.stringWidth(line, c._fontname, size)
            c.drawString(x + max_w - tw, ly, line)
        elif align == 'C':
            tw = c.stringWidth(line, c._fontname, size)
            c.drawString(x + (max_w - tw) / 2, ly, line)
        else:
            c.drawString(x, ly, line)

def _pdf_wrap(c, text, size, bold, italic, max_w):
    """Quebra o texto em linhas que cabem em max_w."""
    font = 'Helvetica-Bold' if bold else ('Helvetica-Oblique' if italic else 'Helvetica')
    words = text.replace('\n', ' \n ').split(' ')
    lines = []
    current = ''
    for word in words:
        if word == '\n':
            lines.append(current.strip())
            current = ''
            continue
        test = (current + ' ' + word).strip() if current else word
        if c.stringWidth(test, font, size) <= max_w:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines

def _pdf_header(c, guide_title=''):
    """Header NF-SP style."""
    _pdf_rect(c, 0, 0, _SW, _H_H, '#A72C31')
    _pdf_txt(c, 0.45, 0.08, 5.5, 0.40, 'SIGMA CONTABILIDADE',
             size=10, bold=True, color='#FFFFFF')
    if guide_title:
        _pdf_txt(c, 6.0, 0.08, 6.88, 0.40, guide_title,
                 size=10, color='#FFFFFF', align='R')

def _pdf_footer(c, page_num, total):
    """Rodapé NF-SP style."""
    footer_text = (f'Sigma Contabilidade  •  Goiânia – GO  •  '
                   f'Dúvidas? Entre em contato: {_SIGMA["tel2"]}')
    _pdf_rect(c, 0, _F_Y, _SW, _F_H, '#F5F5F5')
    _pdf_txt(c, 0.45, 7.16, 11.20, 0.30, footer_text,
             size=8, color='#555555')
    _pdf_txt(c, 12.13, 7.16, 0.80, 0.30,
             f'{page_num}/{total}', size=8, color='#555555', align='R')

def _pdf_label_block(c, x, y, w, h, fill_hex, text, txt_color='#FFFFFF', bold=True):
    _pdf_rect(c, x, y, w, h, fill_hex)
    _pdf_txt(c, x + 0.12, y + 0.04, w - 0.20, h, text,
             size=8, bold=bold, color=txt_color)

def _pdf_content_block(c, x, y, w, h, fill_hex, text, size=11, color='#1A1A1A'):
    _pdf_rect(c, x, y, w, h, fill_hex)
    _pdf_txt(c, x + 0.15, y + 0.10, w - 0.20, h - 0.15,
             text, size=size, color=color)

def _pdf_capa(c, dados):
    """Capa PDF — idêntica ao PPTX."""
    # Fundo branco
    _pdf_rect(c, 0, 0, _SW, _SH, '#FFFFFF')
    # Painel esquerdo
    _pdf_rect(c, 0.0, 0.0, 5.20, _SH, '#A72C31')
    _pdf_rect(c, 5.20, 0.0, 0.06, _SH, '#801E24')

    # SIGMA
    _pdf_txt(c, 0.55, 1.60, 4.20, 1.00, 'SIGMA',
             size=54, bold=True, color='#FFFFFF')
    _pdf_txt(c, 0.55, 2.55, 4.20, 0.50, 'CONTABILIDADE',
             size=18, color='#FFCCCC')
    _pdf_txt(c, 0.55, 3.05, 4.20, 0.40, 'Além da Contabilidade',
             size=12, italic=True, color='#FFDDDD')
    _pdf_rect(c, 0.55, 3.55, 3.50, 0.04, '#FFCCCC')

    # Direita
    subtitulo = dados.get('subtitulo', '') or 'Guia passo a passo — simples e rápido'
    titulo    = dados.get('titulo', 'Informativo Sigma')
    cliente   = dados.get('cliente', '')

    _pdf_txt(c, 5.50, 1.50, 7.40, 0.50, subtitulo,
             size=14, italic=True, color='#A72C31')
    _pdf_txt(c, 5.50, 2.10, 7.40, 2.20, titulo,
             size=36, bold=True, color='#1A1A1A')
    if cliente:
        _pdf_txt(c, 5.50, 4.50, 7.40, 0.50,
                 f'Preparado para: {cliente}', size=13, color='#555555')
    _pdf_txt(c, 0.55, 6.90, 4.20, 0.40,
             f'{_SIGMA["site"]}  •  {_SIGMA["tel1"]}',
             size=10, italic=True, color='#AAAAAA')

def _pdf_passo(c, passo, num_slide, total, titulo_guia, imagem=None):
    """Slide de passo PDF — idêntico ao PPTX."""
    # Fundo branco
    _pdf_rect(c, 0, 0, _SW, _SH, '#FFFFFF')

    _pdf_header(c, titulo_guia)
    _pdf_footer(c, num_slide, total)

    # Barra de passo
    _pdf_rect(c, 0, _B_Y, _SW, _B_H, '#1A3A6E')
    # Badge amarelo
    _pdf_rect(c, _C_X, 0.725, _B_BW, _B_BW, '#FFCD07')
    _pdf_txt(c, _C_X, 0.755, _B_BW, 0.46,
             str(passo.get('num', '')),
             size=20, bold=True, color='#1A3A6E', align='C')
    _pdf_txt(c, 1.13, 0.79, 11.73, 0.46,
             passo.get('titulo', ''), size=17, bold=True, color='#FFFFFF')

    has_img = imagem is not None
    c_w = 6.90 if has_img else _C_W
    c_x = _C_X

    # O QUE FAZER
    y = _C_Y0
    _pdf_label_block(c, c_x, y, c_w, _LBL_H, '#A72C31', 'O QUE FAZER')
    y += _LBL_H
    descricao = passo.get('descricao', '')
    _pdf_content_block(c, c_x, y, c_w, _CONT_H, '#F5F5F5', descricao, size=11)

    # CLIQUE / TOQUE AQUI
    y += _CONT_H
    link = passo.get('link', '')
    _pdf_label_block(c, c_x, y, c_w, _LBL_H, '#1351B4',
                     '\U0001F449  CLIQUE / TOQUE AQUI')
    y += _LBL_H
    link_text = link if link else 'Siga as instruções na tela conforme descrito acima.'
    _pdf_content_block(c, c_x, y, c_w, _LINK_H, '#DBE8FB', link_text, size=11)

    # O QUE VOCÊ VERÁ
    y += _LINK_H
    resultado = passo.get('resultado', '')
    available_h = _F_Y - y - 0.08
    if available_h > 0.4:
        _pdf_label_block(c, c_x, y, c_w, _LBL_H, '#168821',
                         '\U0001f5a5  O QUE VOCÊ VERÁ NA TELA')
        y += _LBL_H
        grn_h = min(available_h - _LBL_H, _GRNH)
        if grn_h > 0.2:
            txt_r = resultado if resultado else 'Confira o resultado na tela e prossiga para o próximo passo.'
            _pdf_content_block(c, c_x, y, c_w, grn_h, '#E3F5E1', txt_r, size=11)

    # Painel direito com imagem
    if has_img:
        _pdf_rect(c, 7.66, _C_Y0, 5.22, 5.62, '#F5F5F5')
        panel_x, panel_y = 7.76, _C_Y0 + 0.15
        panel_w, panel_h = 5.02, 5.32
        img_w, img_h = imagem.size
        ratio = min(panel_w / (img_w / 96), panel_h / (img_h / 96))
        draw_w = (img_w / 96) * ratio
        draw_h = (img_h / 96) * ratio
        off_x = panel_x + (panel_w - draw_w) / 2
        off_y = panel_y + (panel_h - draw_h) / 2
        buf = io.BytesIO()
        imagem.convert('RGB').save(buf, 'JPEG', quality=85)
        buf.seek(0)
        ir = ImageReader(buf)
        rl_x = _pi(off_x)
        rl_y = _PDF_H - _pi(off_y) - _pi(draw_h)
        c.drawImage(ir, rl_x, rl_y, _pi(draw_w), _pi(draw_h))

def _pdf_avisos(c, avisos, num_slide, total, titulo_guia):
    """Slide de avisos PDF."""
    _pdf_rect(c, 0, 0, _SW, _SH, '#FFFFFF')
    _pdf_header(c, titulo_guia)
    _pdf_footer(c, num_slide, total)

    _pdf_rect(c, 0, _B_Y, _SW, _B_H, '#1A3A6E')
    _pdf_txt(c, 0.45, 0.75, 12.43, 0.45,
             '⚠️  ANTES DE COMEÇAR — INFORMAÇÕES IMPORTANTES',
             size=16, bold=True, color='#FFFFFF')

    y_start = _C_Y0
    col_w   = 6.00
    col_gap = 0.43
    box_h   = 1.50
    row_gap = 0.08
    fills   = ['#FFF8DC', '#DBE8FB', '#E3F5E1', '#FFF8DC', '#DBE8FB', '#E3F5E1']
    txcols  = ['#A72C31', '#1A3A6E', '#168821', '#A72C31', '#1A3A6E', '#168821']

    for i, aviso in enumerate(avisos[:6]):
        col = i % 2
        row = i // 2
        x = _C_X + col * (col_w + col_gap)
        y = y_start + row * (box_h + row_gap)
        _pdf_rect(c, x, y, col_w, box_h, fills[i])
        _pdf_txt(c, x + 0.15, y + 0.15, col_w - 0.25, box_h - 0.25,
                 aviso, size=13, bold=True, color=txcols[i])

def _pdf_links(c, links, num_slide, total, titulo_guia):
    """Slide de links PDF."""
    _pdf_rect(c, 0, 0, _SW, _SH, '#FFFFFF')
    _pdf_header(c, titulo_guia)
    _pdf_footer(c, num_slide, total)

    _pdf_rect(c, 0, _B_Y, _SW, _B_H, '#1A3A6E')
    _pdf_txt(c, 0.45, 0.75, 12.43, 0.45,
             '🔗  LINKS E RECURSOS ADICIONAIS',
             size=16, bold=True, color='#FFFFFF')

    y = _C_Y0
    _pdf_rect(c, _C_X, y, _C_W, 0.35, '#A72C31')
    _pdf_txt(c, _C_X + 0.15, y + 0.05, 4.00, 0.25,
             'DESCRIÇÃO', size=9, bold=True, color='#FFFFFF')
    _pdf_txt(c, _C_X + 4.30, y + 0.05, 8.00, 0.25,
             'ENDEREÇO / LINK', size=9, bold=True, color='#FFFFFF')
    y += 0.35

    row_h = 0.55
    for i, lnk in enumerate(links[:9]):
        fill = '#F5F5F5' if i % 2 == 0 else '#FFFFFF'
        _pdf_rect(c, _C_X, y, _C_W, row_h, fill)
        _pdf_txt(c, _C_X + 0.15, y + 0.08, 4.00, row_h - 0.10,
                 lnk.get('label', ''), size=11, color='#1A1A1A')
        _pdf_txt(c, _C_X + 4.30, y + 0.08, 8.00, row_h - 0.10,
                 lnk.get('url', ''), size=10, color='#1351B4')
        y += row_h

def _pdf_contato(c, num_slide, total):
    """Slide de contato final PDF — idêntico ao PPTX (ISSNET style)."""
    _pdf_rect(c, 0, 0, _SW, _SH, '#FFFFFF')

    # Header alto
    _pdf_rect(c, 0.0, 0.0, _SW, _K_H_H, '#A72C31')
    _pdf_rect(c, 0.0, _K_H_H, _SW, 0.06, '#801E24')
    _pdf_txt(c, 0.60, 0.28, 12.00, 0.65,
             'PRECISA DE AJUDA?', size=24, bold=True, color='#FFFFFF')
    _pdf_txt(c, 0.60, 0.82, 12.00, 0.35,
             'Nossa equipe está pronta para te ajudar',
             size=13, italic=True, color='#FFDDDD')

    # Descrição
    _pdf_rect(c, 0.50, _K_DSC_Y, 12.30, _K_DSC_H, '#FDF0F0')
    _pdf_txt(c, 0.75, 1.58, 11.80, 0.95,
             (f'Entre em contato com a Sigma Contabilidade para tirar dúvidas, '
              f'solicitar suporte ou agendar uma conversa com nossos especialistas. '
              f'Estamos disponíveis de {_SIGMA["horario"]}.'),
             size=14, color='#1A1A1A')

    # 4 Cards
    card_data = [
        ('📞', 'Telefones',
         f'{_SIGMA["tel1"]}\n{_SIGMA["tel1_d"]}\n\n'
         f'{_SIGMA["tel2"]}\n{_SIGMA["tel2_d"]}'),
        ('✉️', 'E-mail',
         f'{_SIGMA["email1"]}\n{_SIGMA["email2"]}\n{_SIGMA["email3"]}'),
        ('🌐', 'Site',
         f'{_SIGMA["site"]}\n\nAcesse e conheça\nnossas soluções'),
        ('📍', 'Endereço',
         f'{_SIGMA["end1"]}\n{_SIGMA["end2"]}\n{_SIGMA["end3"]}'),
    ]

    for i, (icon, label, content) in enumerate(card_data):
        cx = _K_CARDS_X[i]
        _pdf_rect(c, cx, _K_CRD_Y, _K_CRD_W, _K_CRD_H, '#F5F5F5')
        _pdf_txt(c, cx + 0.10, _K_CRD_Y + 0.10, 0.60, 0.60,
                 icon, size=22, color='#1A1A1A', align='C')
        _pdf_txt(c, cx + 0.80, _K_CRD_Y + 0.15, 1.80, 0.40,
                 label, size=13, bold=True, color='#A72C31')
        _pdf_txt(c, cx + 0.15, _K_CRD_Y + 0.62, _K_CRD_W - 0.20, 1.70,
                 content, size=11, color='#555555')

    # Rodapé bordeaux
    _pdf_rect(c, 0.0, _K_FT_Y, _SW, _K_FT_H, '#A72C31')
    _pdf_txt(c, 0.60, 5.75, 8.00, 0.65,
             'SIGMA CONTABILIDADE', size=26, bold=True, color='#FFFFFF')
    _pdf_txt(c, 0.60, 6.35, 7.00, 0.40,
             'Além da Contabilidade', size=14, italic=True, color='#FFCCCC')
    _pdf_txt(c, 0.60, 6.85, 12.00, 0.35,
             f'{_SIGMA["site"]}  •  {_SIGMA["tel1"]}  •  {_SIGMA["tel2"]}  •  {_SIGMA["email1"]}',
             size=10, color='#FFCCCC')
    _pdf_txt(c, 12.80, 7.20, 0.45, 0.25,
             f'{num_slide}/{total}', size=9, color='#AAAAAA', align='R')

# ─────────────────────────────────────────────────────────────────────────────
# PDF — GERADOR PRINCIPAL
# ─────────────────────────────────────────────────────────────────────────────

def _gerar_pdf(dados, path_out):
    """Gera o PDF completo e salva em path_out."""
    from reportlab.pdfgen.canvas import Canvas
    c = Canvas(path_out, pagesize=(_PDF_W, _PDF_H))
    c.setTitle(dados.get('titulo', 'Informativo Sigma'))
    c.setAuthor('Sigma Contabilidade')

    titulo_guia = dados.get('titulo', 'Informativo Sigma')
    passos      = dados.get('passos', [])
    avisos      = dados.get('avisos', [])
    links       = dados.get('links', [])
    imagens     = dados.get('imagens', [])

    total = 1
    if avisos: total += 1
    total += len(passos)
    if links: total += 1
    total += 1

    num = 1

    # 1. Capa
    _pdf_capa(c, dados)
    c.showPage()
    num += 1

    # 2. Avisos
    if avisos:
        _pdf_avisos(c, avisos, num, total, titulo_guia)
        c.showPage()
        num += 1

    # 3. Passos
    for i, passo in enumerate(passos):
        img = imagens[i] if i < len(imagens) else None
        _pdf_passo(c, passo, num, total, titulo_guia, img)
        c.showPage()
        num += 1

    # 4. Links
    if links:
        _pdf_links(c, links, num, total, titulo_guia)
        c.showPage()
        num += 1

    # 5. Contato
    _pdf_contato(c, num, total)
    c.showPage()

    c.save()

# ─────────────────────────────────────────────────────────────────────────────
# CLAUDE API — Processar texto orgânico
# ─────────────────────────────────────────────────────────────────────────────

def _processar_texto_claude(texto: str, titulo: str) -> list[dict]:
    """
    Envia o texto orgânico ao Claude e retorna lista de passos estruturados.
    Cada passo: {num, titulo, descricao, resultado, link}
    """
    import anthropic
    api_key = os.getenv('ANTHROPIC_API_KEY', '')
    if not api_key:
        return _fallback_passos(texto)

    prompt = f"""Você é um redator especializado da Sigma Contabilidade.

Recebi o seguinte texto descrevendo um procedimento para clientes:
---
{texto}
---

Título do informativo: {titulo}

Sua tarefa: transformar esse texto em PASSOS NUMERADOS no formato JSON.

Retorne APENAS um JSON válido, sem markdown, sem comentários:
{{
  "passos": [
    {{
      "num": 1,
      "titulo": "Título curto do passo (max 60 chars)",
      "descricao": "Descrição clara do que fazer (2-4 frases)",
      "resultado": "O que o cliente verá/receberá após este passo (1-2 frases)",
      "link": ""
    }}
  ]
}}

Regras:
- Máximo 8 passos
- Cada título começa com verbo no imperativo (Acesse, Clique, Preencha...)
- Descrição: clara, sem termos técnicos
- Resultado: o que aparece na tela ou o que acontece
- link: preencha apenas se o texto mencionar uma URL específica, senão deixe ""
- Se o texto for muito simples, agrupe em 3-5 passos"""

    try:
        client = anthropic.Anthropic(api_key=api_key)
        msg = client.messages.create(
            model='claude-haiku-4-5-20251001',
            max_tokens=2000,
            messages=[{'role': 'user', 'content': prompt}]
        )
        raw = msg.content[0].text.strip()
        # Limpar markdown se veio
        if raw.startswith('```'):
            raw = raw.split('```')[1]
            if raw.startswith('json'):
                raw = raw[4:]
        data = json.loads(raw.strip())
        return data.get('passos', [])
    except Exception as e:
        print(f'[informativos] Claude API erro: {e}')
        return _fallback_passos(texto)

def _fallback_passos(texto: str) -> list[dict]:
    """Divide o texto em passos simples quando a API falha."""
    linhas = [l.strip() for l in texto.split('\n') if l.strip()]
    passos = []
    for i, linha in enumerate(linhas[:6], 1):
        passos.append({
            'num': i,
            'titulo': f'Passo {i}',
            'descricao': linha,
            'resultado': '',
            'link': '',
        })
    return passos

# ─────────────────────────────────────────────────────────────────────────────
# FLASK — ROTAS
# ─────────────────────────────────────────────────────────────────────────────

@informativos_bp.route('/')
def index():
    if login_obrigatorio():
        from flask import redirect, url_for
        return redirect(url_for('auth.login'))
    return render_template('informativos/form.html')

@informativos_bp.route('/gerar', methods=['POST'])
def gerar():
    if login_obrigatorio():
        from flask import redirect, url_for
        return redirect(url_for('auth.login'))

    titulo    = request.form.get('titulo', '').strip()
    subtitulo = request.form.get('subtitulo', '').strip()
    cliente   = request.form.get('cliente', '').strip()
    texto_org = request.form.get('texto_organico', '').strip()

    if not titulo or not texto_org:
        return render_template('informativos/form.html',
                               erro='Preencha o título e o texto do informativo.')

    # Avisos
    avisos = [v.strip() for v in request.form.getlist('aviso') if v.strip()]

    # Links
    links = []
    idx = 1
    while True:
        lbl = request.form.get(f'link_label___{idx}__', '').strip()
        url = request.form.get(f'link_url___{idx}__', '').strip()
        if not lbl and not url:
            break
        if lbl or url:
            links.append({'label': lbl or url, 'url': url})
        idx += 1

    # Imagens (leitura dos uploads)
    imagens = []
    for i in range(1, 30):
        f = request.files.get(f'imagem_{i}')
        if f and f.filename:
            ext = f.filename.rsplit('.', 1)[-1].lower()
            if ext in ALLOWED_IMG:
                try:
                    img = PILImage.open(f.stream)
                    img.load()
                    imagens.append(img.copy())
                except Exception:
                    pass

    # Processar texto com Claude
    passos = _processar_texto_claude(texto_org, titulo)
    if not passos:
        return render_template('informativos/form.html',
                               erro='Não foi possível estruturar o informativo. Tente novamente.')

    # Montar dados completos
    dados = {
        'titulo':    titulo,
        'subtitulo': subtitulo,
        'cliente':   cliente,
        'passos':    passos,
        'avisos':    avisos,
        'links':     links,
        'imagens':   imagens,
    }

    # Gerar arquivos temporários
    token = str(uuid.uuid4())
    tmp_dir = '/tmp'

    pptx_path = os.path.join(tmp_dir, f'{token}.pptx')
    pdf_path  = os.path.join(tmp_dir, f'{token}.pdf')

    try:
        _gerar_pptx(dados, pptx_path)
    except Exception as e:
        print(f'[informativos] PPTX erro: {e}')
        import traceback; traceback.print_exc()
        return render_template('informativos/form.html',
                               erro=f'Erro ao gerar PPTX: {str(e)}')

    try:
        _gerar_pdf(dados, pdf_path)
    except Exception as e:
        print(f'[informativos] PDF erro: {e}')
        import traceback; traceback.print_exc()
        pdf_path = None

    nome_base = titulo[:50].replace(' ', '_').replace('/', '-')

    _CACHE[token] = {
        'pptx': pptx_path,
        'pdf':  pdf_path,
        'nome': nome_base,
        'ts':   time.time(),
    }

    # Limpar cache antigo (> 2h)
    agora = time.time()
    for k in list(_CACHE.keys()):
        if agora - _CACHE[k]['ts'] > 7200:
            for ext in ('pptx', 'pdf'):
                p = _CACHE[k].get(ext)
                if p and os.path.exists(p):
                    try: os.unlink(p)
                    except: pass
            del _CACHE[k]

    return render_template('informativos/resultado.html',
                           token=token,
                           nome=nome_base,
                           titulo=titulo,
                           n_passos=len(passos),
                           pdf_ok=(pdf_path is not None))

@informativos_bp.route('/download/<token>/<tipo>')
def download(token, tipo):
    if login_obrigatorio():
        from flask import redirect, url_for
        return redirect(url_for('auth.login'))

    entrada = _CACHE.get(token)
    if not entrada:
        return render_template('informativos/form.html',
                               erro='Documento expirado. Gere novamente.')

    nome = entrada['nome']

    if tipo == 'pptx':
        path = entrada.get('pptx')
        if not path or not os.path.exists(path):
            return render_template('informativos/form.html',
                                   erro='PPTX não disponível.')
        return send_file(path, as_attachment=True,
                         download_name=f'{nome}.pptx',
                         mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    elif tipo == 'pdf':
        path = entrada.get('pdf')
        if not path or not os.path.exists(path):
            return render_template('informativos/form.html',
                                   erro='PDF não disponível.')
        return send_file(path, as_attachment=True,
                         download_name=f'{nome}.pdf',
                         mimetype='application/pdf')

    return render_template('informativos/form.html', erro='Tipo inválido.')
